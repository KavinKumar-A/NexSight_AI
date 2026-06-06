"""
Generates the NexSight AI hackathon presentation as both PPTX and PDF.
Run: python submission/build_presentation.py
"""
from pathlib import Path

OUT = Path(__file__).resolve().parent
AZURE   = (0x00, 0x78, 0xD4)
DARK    = (0x10, 0x32, 0x55)
INK     = (0x20, 0x20, 0x20)
GREY    = (0x60, 0x60, 0x60)
LIGHT   = (0xF3, 0xF7, 0xFB)
GREEN   = (0x00, 0xC8, 0x78)
RED     = (0xEF, 0x44, 0x44)

# ── Slide content (title, bullets) ─────────────────────────────
SLIDES = [
    {
        "title": "NexSight AI",
        "subtitle": "Manufacturing Intelligence Platform",
        "meta": [
            "Microsoft Build AI Hackathon 2025",
            "Theme: AI Meets Data — From Noise to Insight",
            "Repository: github.com/KavinKumar-A/NexSight_AI",
        ],
        "cover": True,
    },
    {
        "title": "The Problem",
        "bullets": [
            ("Manufacturing lines generate a flood of noisy telemetry", "Thousands of sensor readings per hour across temperature, vibration, humidity, speed and power — most of it never analysed."),
            ("Defects are expensive and discovered too late", "PCB faults (opens, shorts, mousebites...) slip through manual inspection; root causes stay hidden in the data."),
            ("Insight is buried under raw numbers", "Operators see dashboards of metrics, not answers. Quality patterns, failure risks and corrective actions go unnoticed."),
            ("Goal", "Turn raw, noisy shop-floor data into clear, explainable, actionable manufacturing intelligence."),
        ],
    },
    {
        "title": "Our Solution — NexSight AI",
        "bullets": [
            ("One platform, noise to insight", "Ingests sensor telemetry + real PCB inspection imagery and surfaces what actually matters."),
            ("Discovers hidden patterns", "Statistical correlation mining across shifts, machines and product lines (validated p<0.01)."),
            ("Explains root causes", "Feature-importance / SHAP decomposition of every quality driver."),
            ("Predicts and prevents", "XGBoost forecasts defects, yield and machine-failure risk before they happen."),
            ("Sees defects", "Computer-vision PCB defect detection with Grad-CAM visual explanations."),
            ("Acts in real time", "Live WebSocket sensor stream + SSE alerts + an AI assistant for natural-language Q&A."),
        ],
    },
    {
        "title": "Theme Alignment — From Noise to Insight",
        "bullets": [
            ("NOISE  →  raw inputs", "10,000 synthetic telemetry records + 1,500 real DeepPCB image pairs (10,013 annotations)."),
            ("SIGNAL  →  AI engines", "Pattern discovery, anomaly detection, root-cause, prediction and computer vision."),
            ("INSIGHT  →  outputs", "Health score, ranked recommendations, failure forecasts and live alerts."),
            ("ACTION  →  decisions", "Prioritised, ROI-ranked corrective actions an operator can execute today."),
        ],
    },
    {
        "title": "Key Features",
        "bullets": [
            ("Real-time streaming", "WebSocket sensor feed + Server-Sent-Event alert stream."),
            ("Pattern discovery engine", "Surfaces 7+ hidden quality correlations automatically."),
            ("Root-cause intelligence", "SHAP / permutation-importance explainability."),
            ("Predictive analytics", "XGBoost models for defect, yield and failure risk."),
            ("Computer-vision inspection", "DeepPCB defect detection + Grad-CAM heatmaps."),
            ("Health score + anomaly detection", "Composite KPI and Isolation-Forest anomaly flags."),
            ("AI recommendations + assistant", "Ranked actions and natural-language data Q&A."),
        ],
    },
    {
        "title": "Architecture & Tech Stack",
        "bullets": [
            ("Backend", "FastAPI (Python 3.12) — single service exposing REST, WebSocket and SSE."),
            ("Machine learning", "scikit-learn, XGBoost, LightGBM; SHAP for explainability."),
            ("Computer vision", "OpenCV + Pillow with Grad-CAM; works with or without PyTorch."),
            ("Frontend", "Zero-build vanilla JS + Plotly.js + Canvas API single-page dashboard."),
            ("Data layer", "Synthetic telemetry generator + real DeepPCB dataset + cached analytics."),
            ("Delivery", "`pip install -r requirements.txt` then `python run.py` — one command to run."),
        ],
    },
    {
        "title": "AI / ML Approach",
        "bullets": [
            ("Pattern discovery", "Correlation + statistical testing reveals shift / machine / environment effects."),
            ("Root cause", "Model feature importance and SHAP attribute defects to concrete drivers."),
            ("Prediction", "XGBoost regression/classification forecasts yield and failure risk (R² ~0.87)."),
            ("Computer vision", "Defect localisation from DeepPCB annotations; Grad-CAM shows model focus."),
            ("Anomaly detection", "Isolation Forest + domain thresholds flag excursions in real time."),
            ("Explainability first", "Every insight ships with a confidence score and the 'why' behind it."),
        ],
    },
    {
        "title": "Data Strategy",
        "bullets": [
            ("Real data — DeepPCB", "1,500 PCB image pairs, 10,013 defect annotations across 6 defect classes."),
            ("Synthetic data — engineered", "10,000 telemetry records with deliberately embedded quality patterns."),
            ("Why both", "Real data grounds the CV; synthetic data lets us prove pattern-discovery works."),
            ("Live data", "Each session adds streaming WebSocket readings, mixed into the metrics."),
            ("Reproducible", "Data and trained models are committed; caches pre-warm on startup."),
        ],
    },
    {
        "title": "Impact & Results",
        "bullets": [
            ("1,535 defects detected", "Across 200 PCB inspection images with full type + severity breakdown."),
            ("7 hidden patterns surfaced", "Validated with p<0.01 statistical significance."),
            ("Failure risk forecasting", "Flags high-risk machines 48-72 hours ahead for preventive action."),
            ("Quantified savings", "Recommendation engine projects measurable yield gain and cost reduction."),
            ("Sub-2s alerts", "Threshold breaches surface in real time over WebSocket / SSE."),
        ],
    },
    {
        "title": "Live Demo",
        "bullets": [
            ("Executive dashboard", "Health score, KPIs and live machine status at a glance."),
            ("Defect intelligence", "CV results, severity breakdown and defect encyclopedia."),
            ("Patterns & predictions", "Discovered correlations and forecast charts."),
            ("AI assistant", "Ask questions in plain English and get explained answers."),
            ("How to run", "pip install -r requirements.txt  →  python run.py  →  http://localhost:8000"),
        ],
    },
    {
        "title": "Roadmap",
        "bullets": [
            ("Train the CNN end-to-end", "Replace annotation-driven CV with a trained PCBDefectNet + full Grad-CAM."),
            ("Azure integration", "Deploy on Azure App Service; wire Azure OpenAI into the assistant."),
            ("Streaming ingestion", "Connect real PLC / IoT Hub sensor feeds in place of the simulator."),
            ("Automated actioning", "Push recommendations directly to maintenance / work-order systems."),
        ],
    },
    {
        "title": "Thank You",
        "subtitle": "NexSight AI — From Noise to Insight",
        "meta": [
            "Repository: github.com/KavinKumar-A/NexSight_AI",
            "Microsoft Build AI Hackathon 2025",
            "Theme: AI Meets Data — From Noise to Insight",
        ],
        "cover": True,
    },
]


# ════════════════════════════════════════════════════════════════
# PPTX
# ════════════════════════════════════════════════════════════════
def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def rgb(t): return RGBColor(*t)

    def add_rect(slide, x, y, w, h, color):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid(); shp.fill.fore_color.rgb = rgb(color)
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    for s in SLIDES:
        slide = prs.slides.add_slide(blank)
        if s.get("cover"):
            add_rect(slide, 0, 0, SW, SH, DARK)
            add_rect(slide, 0, Inches(3.05), SW, Inches(0.06), AZURE)
            tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.2)).text_frame
            tb.word_wrap = True
            p = tb.paragraphs[0]; r = p.add_run(); r.text = s["title"]
            r.font.size = Pt(54); r.font.bold = True; r.font.color.rgb = rgb((255,255,255))
            if s.get("subtitle"):
                sb = slide.shapes.add_textbox(Inches(0.9), Inches(3.25), Inches(11.5), Inches(0.8)).text_frame
                sb.word_wrap = True
                p = sb.paragraphs[0]; r = p.add_run(); r.text = s["subtitle"]
                r.font.size = Pt(24); r.font.color.rgb = rgb((0x9C,0xD3,0xFF))
            mb = slide.shapes.add_textbox(Inches(0.9), Inches(4.3), Inches(11.5), Inches(2.0)).text_frame
            mb.word_wrap = True
            for i, m in enumerate(s.get("meta", [])):
                p = mb.paragraphs[0] if i == 0 else mb.add_paragraph()
                r = p.add_run(); r.text = m
                r.font.size = Pt(16); r.font.color.rgb = rgb((0xC8,0xD8,0xE8))
                p.space_after = Pt(6)
            continue

        # Content slide
        add_rect(slide, 0, 0, SW, Inches(1.15), AZURE)
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.22), Inches(12), Inches(0.8)).text_frame
        tb.word_wrap = True
        p = tb.paragraphs[0]; r = p.add_run(); r.text = s["title"]
        r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = rgb((255,255,255))

        body = slide.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(11.8), Inches(5.7)).text_frame
        body.word_wrap = True
        bullets = s["bullets"]
        for i, (head, desc) in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            r = p.add_run(); r.text = "▸  " + head
            r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = rgb(DARK)
            p.space_after = Pt(2)
            p2 = body.add_paragraph()
            r2 = p2.add_run(); r2.text = "      " + desc
            r2.font.size = Pt(13); r2.font.color.rgb = rgb(GREY)
            p2.space_after = Pt(10)

    path = OUT / "NexSight_AI_Presentation.pptx"
    prs.save(str(path))
    return path


# ════════════════════════════════════════════════════════════════
# PDF
# ════════════════════════════════════════════════════════════════
def build_pdf():
    from reportlab.lib.pagesizes import landscape
    from reportlab.pdfgen import canvas
    from reportlab.lib.units import inch

    PAGE = (13.333 * inch, 7.5 * inch)
    W, H = PAGE
    path = OUT / "NexSight_AI_Presentation.pdf"
    c = canvas.Canvas(str(path), pagesize=PAGE)

    def col(t): return (t[0]/255, t[1]/255, t[2]/255)

    def wrap(text, font, size, max_w):
        c.setFont(font, size)
        words, lines, cur = text.split(), [], ""
        for w in words:
            trial = (cur + " " + w).strip()
            if c.stringWidth(trial, font, size) <= max_w:
                cur = trial
            else:
                if cur: lines.append(cur)
                cur = w
        if cur: lines.append(cur)
        return lines

    for s in SLIDES:
        if s.get("cover"):
            c.setFillColorRGB(*col(DARK)); c.rect(0, 0, W, H, fill=1, stroke=0)
            c.setFillColorRGB(*col(AZURE)); c.rect(0, H-3.4*inch, W, 0.06*inch, fill=1, stroke=0)
            c.setFillColorRGB(1,1,1); c.setFont("Helvetica-Bold", 50)
            c.drawString(0.9*inch, H-2.7*inch, s["title"])
            if s.get("subtitle"):
                c.setFillColorRGB(*col((0x9C,0xD3,0xFF))); c.setFont("Helvetica", 22)
                c.drawString(0.9*inch, H-3.3*inch, s["subtitle"])
            c.setFillColorRGB(*col((0xC8,0xD8,0xE8))); c.setFont("Helvetica", 14)
            y = H-4.2*inch
            for m in s.get("meta", []):
                c.drawString(0.9*inch, y, m); y -= 0.34*inch
            c.showPage(); continue

        # content
        c.setFillColorRGB(1,1,1); c.rect(0,0,W,H, fill=1, stroke=0)
        c.setFillColorRGB(*col(AZURE)); c.rect(0, H-1.15*inch, W, 1.15*inch, fill=1, stroke=0)
        c.setFillColorRGB(1,1,1); c.setFont("Helvetica-Bold", 28)
        c.drawString(0.7*inch, H-0.78*inch, s["title"])

        y = H - 1.6*inch
        for head, desc in s["bullets"]:
            c.setFillColorRGB(*col(DARK)); c.setFont("Helvetica-Bold", 16)
            c.drawString(0.8*inch, y, "•  " + head)
            y -= 0.28*inch
            c.setFillColorRGB(*col(GREY));
            for line in wrap(desc, "Helvetica", 12, W-2.0*inch):
                c.drawString(1.05*inch, y, line); y -= 0.24*inch
            y -= 0.12*inch
        c.showPage()

    c.save()
    return path


if __name__ == "__main__":
    p1 = build_pptx()
    print(f"[OK] PPTX -> {p1}")
    p2 = build_pdf()
    print(f"[OK] PDF  -> {p2}")
